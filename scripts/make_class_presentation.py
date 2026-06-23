from __future__ import annotations

import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/private/tmp/matplotlib")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs" / "presentation"
OUT_DIR.mkdir(parents=True, exist_ok=True)
DECK_PATH = ROOT / "outputs" / "museum_visitation_end_to_end_presentation.pptx"

NAVY = RGBColor(21, 35, 55)
TEAL = RGBColor(0, 123, 135)
GOLD = RGBColor(235, 167, 52)
MIST = RGBColor(241, 246, 247)
INK = RGBColor(35, 43, 52)
MUTED = RGBColor(93, 104, 116)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 219, 224)


def text_box(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return shape


def bullets(slide, items, x, y, w, h, size=18, color=INK, gap=0):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return shape


def title(slide, heading, subheading=None):
    text_box(slide, heading, 0.55, 0.35, 11.7, 0.55, 25, NAVY, True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(0.98), Inches(1.2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD
    if subheading:
        text_box(slide, subheading, 0.55, 1.08, 11.6, 0.35, 12.5, MUTED)


def footer(slide, n):
    text_box(slide, f"Museum Visitation Forecasting | {n}", 10.35, 7.18, 2.3, 0.22, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=MIST, line=LINE):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    return s


def add_image_fit(slide, path, x, y, w, h):
    if not Path(path).exists():
        card(slide, x, y, w, h)
        text_box(slide, f"Missing image:\n{Path(path).name}", x + 0.15, y + 0.2, w - 0.3, h - 0.4, 14, MUTED)
        return None
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table(slide, df, x, y, w, h, font=11, header_fill=NAVY):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(font)
    for r in range(df.shape[0]):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = str(df.iloc[r, c])
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else RGBColor(248, 250, 251)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = INK
                p.font.size = Pt(font)
    return table


def build_supporting_charts():
    df = pd.read_csv(ROOT / "data/processed/model_dataset.csv", parse_dates=["month"])
    models = pd.read_csv(ROOT / "outputs/tables/model_comparison_table.csv")

    plt.figure(figsize=(11, 4.2))
    plt.plot(df["month"], df["visitors"], color="#007b87", linewidth=2.4)
    plt.fill_between(df["month"], df["visitors"], color="#007b87", alpha=0.12)
    plt.axvspan(pd.Timestamp("2020-03-01"), pd.Timestamp("2021-06-01"), color="#eba734", alpha=0.18, label="COVID disruption window")
    plt.title("Monthly Avila Adobe Visitors", loc="left", fontsize=16, weight="bold")
    plt.ylabel("Visitors")
    plt.xlabel("")
    plt.grid(axis="y", alpha=0.25)
    plt.legend(frameon=False, loc="upper right")
    plt.tight_layout()
    plt.savefig(OUT_DIR / "monthly_visitors.png", dpi=180)
    plt.close()

    plot_df = models.sort_values("R2", ascending=True)
    colors = ["#8d99a6" if x < 0.7 else "#007b87" for x in plot_df["R2"]]
    plt.figure(figsize=(9, 4.6))
    plt.barh(plot_df["Model"], plot_df["R2"], color=colors)
    plt.axvline(0, color="#232b34", linewidth=0.8)
    plt.title("Model Comparison by R-squared", loc="left", fontsize=16, weight="bold")
    plt.xlabel("R-squared")
    plt.grid(axis="x", alpha=0.25)
    plt.tight_layout()
    plt.savefig(OUT_DIR / "model_r2_bar.png", dpi=180)
    plt.close()

    return df, models


def main():
    df, models = build_supporting_charts()
    top_features = pd.read_csv(ROOT / "outputs/tables/random_forest_feature_importance.csv").head(8)
    top_features["importance"] = top_features["importance"].map(lambda x: f"{x:.3f}")
    top_keywords = pd.read_csv(ROOT / "data/processed/google-trends-marketing-v2/keyword_opportunity_ranking.csv").head(8)
    top_keywords = top_keywords[["term", "theme", "opportunity_score"]]
    top_keywords["opportunity_score"] = top_keywords["opportunity_score"].map(lambda x: f"{x:.2f}")

    summary = {
        "months": len(df),
        "start": df["month"].min().strftime("%b %Y"),
        "end": df["month"].max().strftime("%b %Y"),
        "features": df.shape[1] - 2,
        "mean": f"{df['visitors'].mean():,.0f}",
        "max": f"{df['visitors'].max():,.0f}",
        "min": f"{df['visitors'].min():,.0f}",
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_no = 1

    # 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    text_box(s, "Museum Visitation Forecasting", 0.75, 1.0, 11.8, 0.8, 36, WHITE, True)
    text_box(s, "An end-to-end machine learning project using attendance, weather, Google Trends, and calendar seasonality", 0.8, 1.85, 11.2, 0.6, 18, RGBColor(219, 232, 235))
    text_box(s, "SEIS 763 Machine Learning | Avila Adobe / Los Angeles heritage tourism", 0.8, 5.95, 9.8, 0.4, 16, RGBColor(219, 232, 235))
    card(s, 0.8, 3.0, 3.1, 1.35, RGBColor(245, 248, 248), RGBColor(245, 248, 248))
    text_box(s, "Forecasting Question", 1.05, 3.18, 2.5, 0.3, 14, TEAL, True)
    text_box(s, "Can monthly museum attendance be predicted using seasonal, weather, and public-interest signals?", 1.05, 3.58, 2.45, 0.5, 14, INK)
    footer(s, slide_no); slide_no += 1

    # 2
    s = prs.slides.add_slide(blank)
    title(s, "Executive Summary", "The project predicts monthly museum visitors and translates search-demand patterns into practical marketing recommendations.")
    bullets(s, [
        "Built a month-level forecasting dataset covering Jan 2015 to Jun 2021.",
        "Merged museum visitors with weather, Google Trends, calendar indicators, and lag features.",
        "Best model: LASSO-LARS, R2 = 0.783, RMSE = 2,982 visitors, MAE = 2,259 visitors.",
        "Strongest driver: prior-year same-month visitors, confirming seasonality as the main attendance pattern.",
        "Google Trends adds a marketing lens: searchers are often looking for LA history, Olvera Street, family activities, and tourism discovery terms.",
    ], 0.85, 1.55, 11.5, 4.3, 20, INK, 7)
    footer(s, slide_no); slide_no += 1

    # 3
    s = prs.slides.add_slide(blank)
    title(s, "Business Problem", "Avila Adobe is part of a larger destination ecosystem, so forecasting and demand capture both matter.")
    card(s, 0.7, 1.55, 3.75, 4.65)
    text_box(s, "Operational Need", 0.95, 1.85, 3.2, 0.35, 18, NAVY, True)
    bullets(s, ["Plan staffing and programming", "Understand seasonal peaks", "Prepare for low-attendance periods"], 1.0, 2.35, 3.1, 2.3, 18)
    card(s, 4.8, 1.55, 3.75, 4.65)
    text_box(s, "Marketing Need", 5.05, 1.85, 3.2, 0.35, 18, NAVY, True)
    bullets(s, ["Find search terms that match Avila's identity", "Capture visitors before they arrive downtown", "Position Avila as a destination, not only a pass-through stop"], 5.1, 2.35, 3.0, 2.7, 18)
    card(s, 8.9, 1.55, 3.75, 4.65)
    text_box(s, "Analytics Goal", 9.15, 1.85, 3.2, 0.35, 18, NAVY, True)
    bullets(s, ["Forecast monthly visitors", "Identify useful predictors", "Turn model outputs into action-oriented recommendations"], 9.2, 2.35, 3.1, 2.5, 18)
    footer(s, slide_no); slide_no += 1

    # 4
    s = prs.slides.add_slide(blank)
    title(s, "Forecasting Questions", "These questions define the modeling and interpretation work.")
    bullets(s, [
        "Can we accurately predict monthly visitor counts for Avila Adobe?",
        "How much of visitation is explained by seasonality and prior attendance?",
        "Do weather patterns add meaningful predictive value?",
        "Can Google Trends act as a proxy for demand or visitor intent?",
        "Which search themes look like missed marketing opportunities?",
        "What model is accurate enough to support monthly planning decisions?",
    ], 1.0, 1.65, 11.1, 4.5, 22, INK, 8)
    footer(s, slide_no); slide_no += 1

    # 5
    s = prs.slides.add_slide(blank)
    title(s, "Data Sources", f"Final modeling table: {summary['months']} monthly rows, {summary['features']} engineered predictors, no missing values.")
    data_sources = pd.DataFrame([
        ["Museum attendance", "Monthly visitors", "Target variable"],
        ["Weather", "Temp, precipitation, wind", "Contextual predictor"],
        ["Google Trends", "Search interest 0-100", "Demand / marketing proxy"],
        ["Calendar", "Holidays, seasons, school breaks", "Seasonality signal"],
        ["Lag features", "Lag 1, lag 12, rolling mean", "Forecasting memory"],
    ], columns=["Source", "Examples", "Model Role"])
    add_table(s, data_sources, 0.8, 1.55, 11.8, 3.0, 13)
    bullets(s, [
        f"Coverage: {summary['start']} through {summary['end']}.",
        f"Visitor range: {summary['min']} to {summary['max']} monthly visitors; mean = {summary['mean']}.",
        "Dataset grain is monthly, which keeps all sources aligned on the same time unit.",
    ], 0.95, 5.05, 11.2, 1.2, 17)
    footer(s, slide_no); slide_no += 1

    # 6
    s = prs.slides.add_slide(blank)
    title(s, "End-to-End Pipeline", "The project follows a reproducible data science workflow from raw data to presentation outputs.")
    steps = [
        ("Raw data", "museum, weather,\ntrends, calendar"),
        ("Cleaning", "monthly format,\nconsistent columns"),
        ("Feature engineering", "lags, rolling mean,\nseasonal flags"),
        ("Modeling", "linear, LASSO,\nRF, boosting"),
        ("Evaluation", "RMSE, MAE,\nR2"),
        ("Insights", "forecasting +\nmarketing actions"),
    ]
    x = 0.55
    for i, (h, b) in enumerate(steps):
        card(s, x, 2.25, 1.7, 1.25, WHITE, TEAL)
        text_box(s, h, x + 0.12, 2.42, 1.45, 0.25, 13, NAVY, True, PP_ALIGN.CENTER)
        text_box(s, b, x + 0.12, 2.75, 1.45, 0.45, 11, MUTED, False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            text_box(s, ">", x + 1.78, 2.58, 0.35, 0.3, 22, GOLD, True, PP_ALIGN.CENTER)
        x += 2.05
    bullets(s, [
        "Scripts and notebooks document the path from raw files to processed datasets, model outputs, plots, and presentation visuals.",
        "Google Trends work is versioned so new collection attempts can be compared without overwriting prior runs.",
    ], 1.05, 4.55, 11.0, 1.1, 17)
    footer(s, slide_no); slide_no += 1

    # 7
    s = prs.slides.add_slide(blank)
    title(s, "Visitor Pattern", "Attendance is highly seasonal and sharply disrupted around the COVID period.")
    add_image_fit(s, OUT_DIR / "monthly_visitors.png", 0.8, 1.45, 11.8, 4.6)
    bullets(s, ["This pattern explains why lag and calendar features matter: the model needs memory of prior seasonal cycles."], 1.0, 6.25, 11.0, 0.5, 16)
    footer(s, slide_no); slide_no += 1

    # 8
    s = prs.slides.add_slide(blank)
    title(s, "Feature Engineering", "The final model combines direct signals, seasonal indicators, and time-aware predictors.")
    feature_table = pd.DataFrame([
        ["Weather", "avg_temp_F, min_temp_F, max_temp_F, precipitation, wind"],
        ["Seasonality", "month number, quarter, season flags, month_sin, month_cos"],
        ["Calendar", "federal holidays, school breaks, event season, holiday season"],
        ["Google Trends", "museum-specific, tourism, family, seasonal, cultural heritage terms"],
        ["Lag features", "visitors_lag1, visitors_lag12, rolling_mean_3"],
    ], columns=["Feature Group", "Examples"])
    add_table(s, feature_table, 0.8, 1.55, 11.8, 3.3, 13)
    bullets(s, [
        "Lag 12 captures last year's same-month baseline.",
        "Rolling mean smooths recent momentum.",
        "Trend terms are normalized Google Trends values, not actual search counts.",
    ], 1.0, 5.25, 11.0, 1.0, 17)
    footer(s, slide_no); slide_no += 1

    # 9
    s = prs.slides.add_slide(blank)
    title(s, "Google Trends Research Design", "Search interest is used as both a modeling signal and a marketing opportunity scan.")
    bullets(s, [
        "Primary focus: Avila Adobe and the surrounding El Pueblo / Olvera Street visitor path.",
        "Keyword lanes: museum-specific, discovery marketing, cultural heritage, seasonal intent, family activities, and benchmarks.",
        "Benchmarks: Getty Museum, La Brea Tar Pits, and Griffith Observatory provide context, not the main target.",
        "Interpretation rule: Google Trends values are relative 0-100 interest scores, not raw search volume.",
    ], 0.8, 1.55, 5.7, 4.8, 18, INK, 7)
    add_image_fit(s, ROOT / "outputs/plots/google_trends_marketing/google-trends-marketing-v6/top_keyword_opportunities.png", 6.85, 1.45, 5.75, 4.75)
    footer(s, slide_no); slide_no += 1

    # 10
    s = prs.slides.add_slide(blank)
    title(s, "Top Search Opportunities", "Highest-opportunity keywords suggest how Avila could intercept existing demand.")
    add_table(s, top_keywords.rename(columns={"term": "Term", "theme": "Theme", "opportunity_score": "Score"}), 0.7, 1.45, 6.0, 4.6, 10.5)
    add_image_fit(s, ROOT / "outputs/plots/google_trends_marketing/google-trends-marketing-v6/theme_opportunity_scores.png", 7.05, 1.55, 5.55, 4.45)
    bullets(s, ["Top themes above the opportunity threshold include competitor context, tourism discovery, and lifestyle/history."], 0.85, 6.25, 11.3, 0.5, 16)
    footer(s, slide_no); slide_no += 1

    # 11
    s = prs.slides.add_slide(blank)
    title(s, "Visitor Alignment", "Search terms connected to LA history and Olvera Street align well with the museum visitation story.")
    add_image_fit(s, ROOT / "outputs/plots/google_trends_marketing/google-trends-marketing-v6/avila_vs_top_trends_scaled.png", 0.75, 1.35, 5.85, 4.85)
    add_image_fit(s, ROOT / "outputs/plots/google_trends_marketing/google-trends-marketing-v6/avila_vs_benchmarks_scaled.png", 6.9, 1.35, 5.7, 4.85)
    footer(s, slide_no); slide_no += 1

    # 12
    s = prs.slides.add_slide(blank)
    title(s, "Modeling Approach", "Multiple algorithms were compared using the same final dataset.")
    model_steps = pd.DataFrame([
        ["Baseline", "Linear Regression", "Simple benchmark"],
        ["Regularized", "LASSO, LASSO-LARS", "Handles many correlated predictors"],
        ["Tree ensemble", "Random Forest", "Captures nonlinear relationships"],
        ["Boosting", "Gradient Boosting", "Sequentially improves prediction errors"],
    ], columns=["Category", "Models", "Why Included"])
    add_table(s, model_steps, 0.8, 1.55, 11.8, 2.7, 13)
    bullets(s, [
        "Metrics: RMSE, MAE, and R2.",
        "Interpretation priority: balance predictive accuracy with clear explanation for class presentation.",
        "Regularization is especially useful because many trend and calendar variables can overlap.",
    ], 1.0, 4.65, 11.0, 1.4, 17)
    footer(s, slide_no); slide_no += 1

    # 13
    s = prs.slides.add_slide(blank)
    title(s, "Model Results", "LASSO and LASSO-LARS clearly outperform the baseline and tree models.")
    mt = models.copy()
    mt["RMSE"] = mt["RMSE"].map(lambda x: f"{x:,.0f}")
    mt["MAE"] = mt["MAE"].map(lambda x: f"{x:,.0f}")
    mt["R2"] = mt["R2"].map(lambda x: f"{x:.3f}")
    add_table(s, mt, 0.7, 1.45, 6.1, 3.4, 10.5)
    add_image_fit(s, OUT_DIR / "model_r2_bar.png", 7.05, 1.45, 5.55, 3.7)
    bullets(s, ["Best result: LASSO-LARS with RMSE = 2,982, MAE = 2,259, and R2 = 0.783."], 0.9, 5.7, 11.2, 0.6, 18)
    footer(s, slide_no); slide_no += 1

    # 14
    s = prs.slides.add_slide(blank)
    title(s, "Best Model Fit", "Predicted and actual visitors track the broad monthly pattern closely.")
    add_image_fit(s, ROOT / "outputs/plots/lasso_actual_vs_predicted.png", 0.85, 1.35, 5.75, 4.75)
    add_image_fit(s, ROOT / "outputs/plots/lasso_residuals_over_time.png", 6.9, 1.35, 5.7, 4.75)
    footer(s, slide_no); slide_no += 1

    # 15
    s = prs.slides.add_slide(blank)
    title(s, "Important Predictors", "The model story is dominated by seasonality, with search and weather adding useful context.")
    add_table(s, top_features.rename(columns={"feature": "Feature", "importance": "Importance"}), 0.7, 1.45, 5.9, 4.8, 10.5)
    add_image_fit(s, ROOT / "outputs/plots/lasso_top_predictors_plot.png", 6.9, 1.45, 5.7, 4.8)
    footer(s, slide_no); slide_no += 1

    # 16
    s = prs.slides.add_slide(blank)
    title(s, "Interpretation", "What the model says about visitor behavior.")
    bullets(s, [
        "Seasonality is the backbone: visitors_lag12 is the strongest single predictor in the feature-importance output.",
        "Tourism discovery terms matter: searches for LA attractions and Olvera Street help represent public intent.",
        "Weather can still influence short-term outing decisions, especially precipitation and temperature variables.",
        "A practical forecast should combine last year's same month, recent attendance momentum, and search-interest signals.",
    ], 0.9, 1.55, 11.4, 4.8, 21, INK, 8)
    footer(s, slide_no); slide_no += 1

    # 17
    s = prs.slides.add_slide(blank)
    title(s, "Recommendations", "Use the forecast for planning and the Trends analysis for demand capture.")
    bullets(s, [
        "Operational planning: use monthly forecasts to anticipate staffing, programming, and visitor-service needs.",
        "Marketing: test search copy around 'Los Angeles history,' 'Olvera Street,' 'oldest house in Los Angeles,' and family outing terms.",
        "Seasonal campaigns: align content with winter, rainy-day, school-break, and short-visit search behavior.",
        "Programming: create repeatable hooks such as heritage weekends, family activities, school-break tours, and time-bounded history experiences.",
        "Measurement: track future visitor counts against search-interest changes to evaluate whether campaigns shift demand.",
    ], 0.85, 1.45, 11.6, 5.0, 20, INK, 8)
    footer(s, slide_no); slide_no += 1

    # 18
    s = prs.slides.add_slide(blank)
    title(s, "Limitations And Next Steps", "The model is useful, but the project has clear boundaries.")
    card(s, 0.8, 1.45, 5.7, 4.7)
    text_box(s, "Limitations", 1.05, 1.75, 5.0, 0.35, 19, NAVY, True)
    bullets(s, [
        "Only 78 monthly rows in the final model dataset.",
        "COVID creates an unusual disruption period.",
        "Google Trends values are relative scores, not exact search counts.",
        "Search terms can be sparse or return zeros when interest is too low.",
    ], 1.05, 2.25, 5.0, 3.1, 17)
    card(s, 6.85, 1.45, 5.7, 4.7)
    text_box(s, "Next Steps", 7.1, 1.75, 5.0, 0.35, 19, NAVY, True)
    bullets(s, [
        "Add newer visitor data after June 2021.",
        "Test time-series validation and forecast horizons.",
        "Separate COVID or reopening effects explicitly.",
        "Compare campaign periods against future Trends and visitor changes.",
    ], 7.1, 2.25, 5.0, 3.1, 17)
    footer(s, slide_no); slide_no += 1

    # 19
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    text_box(s, "Conclusion", 0.85, 0.9, 11.8, 0.7, 34, WHITE, True)
    bullets(s, [
        "The project shows that monthly museum visitation can be forecast with useful accuracy.",
        "Seasonality is the strongest driver, but Google Trends helps explain demand and marketing opportunity.",
        "The best presentation takeaway: Avila Adobe can plan around predictable seasonal patterns while using search-interest insights to become a more intentional destination.",
    ], 1.0, 2.05, 11.1, 3.2, 23, RGBColor(231, 240, 242), 9)
    text_box(s, "Questions", 0.95, 6.2, 11.5, 0.5, 24, GOLD, True, PP_ALIGN.CENTER)
    footer(s, slide_no); slide_no += 1

    prs.save(DECK_PATH)
    print(DECK_PATH)


if __name__ == "__main__":
    main()
